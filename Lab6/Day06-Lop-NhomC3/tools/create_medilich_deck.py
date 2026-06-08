from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = "MediLich_DemoDay_Batch03.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

COLORS = {
    "green": RGBColor(27, 107, 90),
    "green2": RGBColor(52, 132, 113),
    "mint": RGBColor(229, 245, 239),
    "mint2": RGBColor(244, 250, 248),
    "navy": RGBColor(28, 41, 51),
    "gray": RGBColor(93, 105, 116),
    "light": RGBColor(246, 248, 250),
    "white": RGBColor(255, 255, 255),
    "amber": RGBColor(244, 176, 55),
    "red": RGBColor(186, 26, 26),
    "blue": RGBColor(42, 107, 177),
    "purple": RGBColor(111, 80, 160),
}


def add_bg(slide, color="white"):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS[color]


def add_text(slide, text, x, y, w, h, size=24, color="navy", bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        add_text(slide, kicker.upper(), 0.65, 0.35, 5.8, 0.3, 10, "green", True)
    add_text(slide, title, 0.65, 0.72, 8.4, 0.62, 27, "navy", True)
    if subtitle:
        add_text(slide, subtitle, 0.67, 1.35, 8.6, 0.35, 13, "gray")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.88), Inches(1.25), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["green"]
    line.line.fill.background()


def add_footer(slide, num):
    add_text(slide, f"{num:02d}", 12.45, 7.05, 0.45, 0.2, 9, "gray", True, PP_ALIGN.RIGHT)
    add_text(slide, "MediLịch · Batch 03 Demo Day", 0.65, 7.05, 3.5, 0.2, 8, "gray")


def add_chip(slide, text, x, y, w, color="green"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]
    return shape


def add_card(slide, x, y, w, h, title, bullets=None, accent="green", body_size=13):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = RGBColor(221, 229, 225)
    card.line.width = Pt(0.8)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.44, 0.35, 15, accent, True)
    if bullets:
        box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.62), Inches(w - 0.45), Inches(h - 0.75))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.name = "Aptos"
            p.font.size = Pt(body_size)
            p.font.color.rgb = COLORS["navy"]
            p.space_after = Pt(5)
    return card


def add_metric(slide, x, y, value, label, color="green"):
    add_text(slide, value, x, y, 1.65, 0.45, 25, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, x - 0.15, y + 0.48, 1.95, 0.34, 9, "gray", False, PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color="gray"):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True
    return line


def add_node(slide, x, y, w, h, title, subtitle="", color="mint"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.color.rgb = COLORS["green"]
    shape.line.width = Pt(1)
    add_text(slide, title, x + 0.12, y + 0.12, w - 0.24, 0.25, 12, "navy", True, PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, x + 0.14, y + 0.42, w - 0.28, 0.28, 8, "gray", False, PP_ALIGN.CENTER)
    return shape


def add_phone(slide, x, y, title, rows, accent="green"):
    shell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.05), Inches(3.75))
    shell.fill.solid()
    shell.fill.fore_color.rgb = RGBColor(19, 30, 38)
    shell.line.fill.background()
    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.12), Inches(y + 0.16), Inches(1.81), Inches(3.43))
    screen.fill.solid()
    screen.fill.fore_color.rgb = COLORS["mint2"]
    screen.line.fill.background()
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.12), Inches(y + 0.16), Inches(1.81), Inches(0.48))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()
    add_text(slide, title, x + 0.25, y + 0.29, 1.5, 0.2, 8, "white", True, PP_ALIGN.CENTER)
    cy = y + 0.83
    for row in rows:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.28), Inches(cy), Inches(1.49), Inches(0.38))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS["white"]
        box.line.color.rgb = RGBColor(219, 227, 224)
        add_text(slide, row, x + 0.36, cy + 0.1, 1.34, 0.14, 6.8, "navy")
        cy += 0.48
    return shell


slides = [
    {
        "title": "MediLịch",
        "subtitle": "Scan đơn thuốc → Lịch uống thông minh",
        "kicker": "Batch 03 · AI Product Demo",
    },
    {
        "title": "Vấn Đề & Tầm Nhìn",
        "subtitle": "Pain không nằm ở nhắc thuốc, mà ở chuyển đơn thành lịch.",
        "kicker": "Problem Statement",
    },
    {
        "title": "SPEC: Yêu Cầu Sản Phẩm",
        "subtitle": "Một lát cắt hẹp, giá trị rõ.",
        "kicker": "Functional & Non-functional",
    },
    {
        "title": "Tech Stack",
        "subtitle": "Stack gọn, đủ nhanh cho MVP chạy thật.",
        "kicker": "Solution Design",
    },
    {
        "title": "UX Research",
        "subtitle": "Người dùng cần ít bước, ít rủi ro.",
        "kicker": "Persona & Journey",
    },
    {
        "title": "Kiến Trúc Hệ Thống",
        "subtitle": "Client nhẹ, backend kiểm soát rủi ro AI.",
        "kicker": "System Architecture",
    },
    {
        "title": "Tính Năng 1: AI Scan Có Kiểm Soát",
        "subtitle": "AI tạo bản nháp, người dùng là người quyết định.",
        "kicker": "Core Feature",
    },
    {
        "title": "Tính Năng 2: Lịch Nhắc + Thẻ Thuốc",
        "subtitle": "Biến đơn thuốc thành hành động hằng ngày.",
        "kicker": "Core Feature",
    },
    {
        "title": "Demo Prototype",
        "subtitle": "Code thật, giao diện thật, flow hoàn chỉnh.",
        "kicker": "Product Evidence",
    },
    {
        "title": "Kết Quả & Metrics",
        "subtitle": "Đo theo readiness demo và giá trị thực tiễn.",
        "kicker": "Validation",
    },
    {
        "title": "Reflection",
        "subtitle": "Thách thức đã vượt qua và bài học rút ra.",
        "kicker": "Learning Agenda",
    },
    {
        "title": "Q&A",
        "subtitle": "MediLịch · Uống đúng, hiểu đúng.",
        "kicker": "Thank You",
    },
]


for idx, meta in enumerate(slides, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "mint2")
    add_title(slide, meta["title"], meta["subtitle"], meta["kicker"])
    add_footer(slide, idx)

    if idx == 1:
        add_text(slide, "MEDILỊCH", 0.85, 2.25, 4.8, 0.65, 42, "green", True)
        add_text(slide, "Scan đơn thuốc → Lịch uống + Thẻ thuốc", 0.9, 3.05, 5.7, 0.35, 18, "navy")
        add_text(slide, "Team Identity", 0.92, 3.72, 1.6, 0.25, 11, "gray", True)
        add_card(slide, 0.85, 4.08, 4.8, 1.65, "Batch 03 · Healthcare AI", [
            "Domain: medication adherence",
            "Approach: human-in-the-loop AI",
            "Outcome: lịch uống an toàn hơn",
        ], "green", 13)
        add_phone(slide, 9.55, 2.0, "Quét đơn", ["Ảnh đơn", "Demo", "Phân tích AI", "Xác nhận"], "green")
        add_phone(slide, 7.2, 2.55, "Nhắc", ["08:00 Amoxicillin", "14:00 Paracetamol", "21:00 Vitamin C"], "green2")
        add_chip(slide, "SPEC", 6.95, 6.0, 0.85)
        add_chip(slide, "UX", 7.95, 6.0, 0.7, "blue")
        add_chip(slide, "PROTOTYPE", 8.8, 6.0, 1.25, "purple")
        add_chip(slide, "REFLECTION", 10.2, 6.0, 1.35, "amber")

    elif idx == 2:
        add_card(slide, 0.8, 2.25, 3.45, 3.25, "Current Pain", [
            "Nhập lịch thủ công",
            "Dễ sai tần suất",
            "Không hiểu tên thuốc",
            "App chưa scan đơn Việt",
        ], "red")
        add_card(slide, 4.75, 2.25, 3.45, 3.25, "Strategic Insight", [
            "OCR chưa đủ giá trị",
            "Cần chuyển thành hành động",
            "Cần kiểm soát rủi ro",
            "User phải xác nhận",
        ], "green")
        add_card(slide, 8.7, 2.25, 3.45, 3.25, "Vision", [
            "Uống đúng giờ",
            "Hiểu đúng thuốc",
            "Giảm nhập tay",
            "An toàn trước tốc độ",
        ], "blue")

    elif idx == 3:
        add_card(slide, 0.78, 2.15, 3.6, 3.95, "Functional Requirements", [
            "Upload/chụp ảnh đơn thuốc",
            "AI trích xuất từng dòng",
            "Review trước khi lưu",
            "Tự tạo lịch nhắc",
            "Hiển thị thẻ thuốc",
            "Tìm nhà thuốc gần",
        ], "green", 12)
        add_card(slide, 4.85, 2.15, 3.6, 3.95, "Non-functional Requirements", [
            "API key không lộ client",
            "Có fixture demo offline",
            "Mobile-first 390px",
            "Không thay tư vấn y tế",
            "Backend trả JSON an toàn",
            "Luồng demo dưới 5 phút",
        ], "blue", 12)
        add_card(slide, 8.92, 2.15, 3.55, 3.95, "Acceptance Criteria", [
            "3 preset chạy được",
            "Chặn lỗi nguy hiểm",
            "Lịch sinh đúng slot",
            "Drug cards render đúng",
            "Citation chỉ khi khớp",
            "Server chạy localhost",
        ], "purple", 12)

    elif idx == 4:
        add_card(slide, 0.8, 2.12, 2.75, 3.8, "Frontend", [
            "HTML, CSS, Vanilla JS",
            "ES Modules",
            "SessionStorage",
            "Mobile Android shell",
        ], "green", 12)
        add_card(slide, 3.85, 2.12, 2.75, 3.8, "Backend", [
            "Node.js + Express",
            "Multer upload memory",
            "dotenv API key",
            "REST API endpoints",
        ], "blue", 12)
        add_card(slide, 6.9, 2.12, 2.75, 3.8, "AI & OCR", [
            "OpenAI Vision",
            "OpenAI drug info",
            "VietOCR optional",
            "JSON normalization",
        ], "purple", 12)
        add_card(slide, 9.95, 2.12, 2.75, 3.8, "Data & Services", [
            "Local drugs.json",
            "Vinmec citations",
            "Overpass nearby",
            "Fixture fallback",
        ], "amber", 12)

    elif idx == 5:
        add_card(slide, 0.8, 2.1, 3.4, 4.2, "Persona", [
            "Người bệnh lớn tuổi",
            "Người nhà chăm sóc",
            "Quen smartphone cơ bản",
            "Sợ nhập sai lịch",
            "Cần chữ rõ",
        ], "green", 12)
        steps = [
            ("Nhận đơn", "Giấy / screenshot"),
            ("Chụp ảnh", "Một thao tác"),
            ("AI đọc", "Bản nháp"),
            ("Review", "Sửa từng dòng"),
            ("Lưu lịch", "Nhắc uống"),
            ("Xem thuốc", "Hiểu rõ hơn"),
        ]
        x = 4.65
        for i, (t, s) in enumerate(steps):
            add_node(slide, x + (i % 3) * 2.35, 2.35 + (i // 3) * 1.62, 1.75, 0.82, t, s, "white")
            if i % 3 != 2:
                add_arrow(slide, x + (i % 3) * 2.35 + 1.78, 2.75 + (i // 3) * 1.62, x + (i % 3 + 1) * 2.35 - 0.08, 2.75 + (i // 3) * 1.62, "green")
        add_text(slide, "Wireframe flow: Scan → Review → Nhắc → Lịch → Thuốc", 4.7, 5.92, 6.2, 0.28, 13, "navy", True)

    elif idx == 6:
        add_node(slide, 0.85, 3.15, 1.75, 0.82, "Browser", "index.html + JS", "white")
        add_node(slide, 3.15, 3.15, 1.9, 0.82, "API Client", "js/api.js", "white")
        add_node(slide, 5.55, 3.15, 1.95, 0.82, "Express", "server/index.js", "white")
        add_node(slide, 8.08, 2.2, 1.95, 0.82, "OpenAI", "Vision + Chat", "mint")
        add_node(slide, 8.08, 4.12, 1.95, 0.82, "VietOCR", "optional sidecar", "mint")
        add_node(slide, 10.55, 2.2, 1.95, 0.82, "Vinmec", "citation lookup", "white")
        add_node(slide, 10.55, 4.12, 1.95, 0.82, "Overpass", "nearby places", "white")
        add_arrow(slide, 2.62, 3.55, 3.12, 3.55, "green")
        add_arrow(slide, 5.07, 3.55, 5.53, 3.55, "green")
        add_arrow(slide, 7.52, 3.38, 8.05, 2.68, "green")
        add_arrow(slide, 7.52, 3.75, 8.05, 4.5, "green")
        add_arrow(slide, 10.05, 2.62, 10.52, 2.62, "green")
        add_arrow(slide, 10.05, 4.55, 10.52, 4.55, "green")
        add_card(slide, 0.85, 5.45, 11.65, 0.75, "Main Flow", [
            "Image → OCR/AI → Review → Schedule → Reminder"
        ], "green", 14)

    elif idx == 7:
        add_card(slide, 0.8, 2.2, 3.25, 3.9, "Business Value", [
            "Giảm nhập tay",
            "Tạo bản nháp nhanh",
            "User xác nhận cuối",
            "Chặn rủi ro cao",
        ], "green", 13)
        add_card(slide, 4.45, 2.2, 3.25, 3.9, "Technical Flow", [
            "Multer nhận ảnh",
            "OpenAI Vision parse",
            "VietOCR khi bật",
            "normalizeLines()",
            "reviewDrugNames()",
        ], "blue", 13)
        add_card(slide, 8.1, 2.2, 3.25, 3.9, "Safety Guardrails", [
            "Confidence badge vàng",
            "Chặn ngoài 1–4",
            "Rule Amoxicillin risky",
            "Disable nút lưu",
        ], "red", 13)

    elif idx == 8:
        add_card(slide, 0.8, 2.18, 3.55, 3.95, "Reminder Engine", [
            "buildSchedule() sinh events",
            "Slot theo tần suất",
            "Taken state trong Set",
            "Toast giả lập notification",
        ], "green", 12)
        add_card(slide, 4.85, 2.18, 3.55, 3.95, "Drug Cards", [
            "matchDrug() local trước",
            "AI cho thuốc thiếu",
            "Client aiCache",
            "Vinmec citation thật",
        ], "purple", 12)
        add_card(slide, 8.9, 2.18, 3.55, 3.95, "Nearby Support", [
            "Geolocation browser",
            "Overpass OSM query",
            "Pharmacy / hospital",
            "Không hứa tồn kho",
        ], "blue", 12)

    elif idx == 9:
        add_phone(slide, 0.95, 2.3, "Scan", ["Ảnh đơn", "Demo", "Tải mẫu", "Phân tích AI"], "green")
        add_phone(slide, 3.35, 2.3, "Review", ["Amoxicillin", "1 viên", "3 lần/ngày", "Lưu & sync"], "blue")
        add_phone(slide, 5.75, 2.3, "Nhắc", ["08:00", "Hôm nay", "Đã uống", "Nhắc lại 10p"], "green2")
        add_phone(slide, 8.15, 2.3, "Lịch", ["Tháng 6/2026", "Ngày có chấm", "Đồng bộ", "Sheet 2 cột"], "purple")
        add_phone(slide, 10.55, 2.3, "Thuốc", ["Thẻ thuốc", "Cách uống", "Vinmec", "Nhà thuốc gần"], "amber")
        add_text(slide, "Demo paths: Happy · Low-confidence · Failure · Correction", 1.15, 6.35, 8.5, 0.28, 13, "navy", True)

    elif idx == 10:
        add_metric(slide, 0.95, 2.35, "3", "preset demo")
        add_metric(slide, 2.85, 2.35, "≤30s", "scan target", "blue")
        add_metric(slide, 4.75, 2.35, "1–4", "frequency guard", "red")
        add_metric(slide, 6.65, 2.35, "390px", "mobile width", "purple")
        add_metric(slide, 8.55, 2.35, "0", "client key leak", "green")
        add_metric(slide, 10.45, 2.35, "5", "core tabs/flows", "amber")
        add_card(slide, 0.9, 4.15, 5.35, 1.45, "Readiness Evidence", [
            "Fixture fallback works",
            "Schedule sinh đúng slot",
            "Drug cards render theo đơn",
        ], "green", 13)
        add_card(slide, 6.75, 4.15, 5.35, 1.45, "Risk Evidence", [
            "Failure path bị chặn",
            "Citation không bịa URL",
            "Backend giữ API key",
        ], "red", 13)

    elif idx == 11:
        add_card(slide, 0.8, 2.15, 3.55, 3.95, "Challenges", [
            "OCR tiếng Việt nhiễu",
            "Sai tần suất nguy hiểm",
            "Citation dễ hallucinate",
            "Demo cần fallback",
            "Scope dễ phình to",
        ], "red", 12)
        add_card(slide, 4.85, 2.15, 3.55, 3.95, "How We Solved", [
            "Augmentation, not automation",
            "Review bắt buộc",
            "Fixture demo offline",
            "Citation lookup thật",
            "Giữ scope hẹp",
        ], "green", 12)
        add_card(slide, 8.9, 2.15, 3.55, 3.95, "Lessons", [
            "AI tốt cho bản nháp",
            "Healthcare cần human loop",
            "Failure path quan trọng",
            "SPEC phải khớp code",
            "Demo cần kịch bản rõ",
        ], "blue", 12)

    elif idx == 12:
        add_text(slide, "MediLịch giúp người bệnh chuyển đơn thuốc thành lịch uống an toàn.", 1.25, 2.38, 10.8, 0.52, 26, "green", True, PP_ALIGN.CENTER)
        add_card(slide, 2.15, 3.45, 3.1, 1.35, "Demo", [
            "Sẵn sàng trình diễn prototype",
            "Happy · Failure · Correction",
        ], "green", 12)
        add_card(slide, 5.55, 3.45, 3.1, 1.35, "Repo", [
            "Day5_group_lab",
            "prototype/server → npm start",
        ], "blue", 12)
        add_card(slide, 8.95, 3.45, 3.1, 1.35, "Contact", [
            "Team MediLịch",
            "Batch 03 Demo Day",
        ], "purple", 12)
        add_text(slide, "Q&A", 5.72, 5.65, 1.8, 0.42, 28, "navy", True, PP_ALIGN.CENTER)


prs.save(OUT)
print(OUT)
